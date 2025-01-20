import pdfplumber
import pandas as pd
from io import BytesIO
from abc import ABC, abstractmethod
from docx import Document
from pptx import Presentation


class BaseExtractor(ABC):
    def __init__(self, file: bytes):
        self.file = file

    @abstractmethod
    def extract(self):
        ...


class PDFExtractor(BaseExtractor):
    def extract(self):
        with pdfplumber.open(BytesIO(self.file)) as pdf:
            text = "".join([page.extract_text() for page in pdf.pages])
            return text.replace("\n", " ")
        

class XLSXExtractor(BaseExtractor):
    # TODO: figure out how to return data properly
    def extract(self):
        data = pd.read_excel(BytesIO(self.file))
        return data.head().to_csv()


class DOCXExtractor(BaseExtractor):
    def extract(self):
        try:
            doc = Document(BytesIO(self.file))
            return " ".join([p.text for p in doc.paragraphs])
        except Exception as e:
            return {"error": f"Failed to extract text: {str(e)}"}
        

class PPTXExtractor(BaseExtractor):
    def extract(self):
        try:
            presentation = Presentation(BytesIO(self.file))
            text = ""
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text += shape.text + " "
            return text
        except Exception as e:
            return {"error": f"Failed to extract text: {str(e)}"}
    

class TXTExtractor(BaseExtractor):
    def extract(self):
        try:
            return self.file.decode('utf-8').replace("\n", " ")
        except Exception as e:
            return {"error": f"Failed to extract text: {str(e)}"}
        

class TextExtractor(BaseExtractor):
    file_types = {
        ".pdf": PDFExtractor,
        ".xlsx": XLSXExtractor,
        ".xls": XLSXExtractor,
        ".docx": DOCXExtractor,
        ".doc": DOCXExtractor,
        ".pptx": PPTXExtractor,
        ".ppt": PPTXExtractor,
        ".txt": TXTExtractor,
    }

    def __init__(self, file: bytes, file_ext: str):
        super().__init__(file)
        self.file_ext = file_ext

    def extract(self):
        file_ext = self.file_ext.lower()
        if file_ext not in self.file_types:
            raise ValueError(f"Unsupported file type: {file_ext}")
        
        extractor_class = self.file_types[file_ext]
        extractor: BaseExtractor = extractor_class(self.file)
        
        return extractor.extract()
